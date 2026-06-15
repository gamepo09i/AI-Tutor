from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT_DIR = Path("slides/weekly")

INK = RGBColor(23, 32, 51)
MUTED = RGBColor(86, 101, 123)
WHITE = RGBColor(255, 255, 255)
LIGHT = RGBColor(244, 247, 251)
LINE = RGBColor(217, 226, 238)
DARK = RGBColor(15, 23, 42)
BLUE = RGBColor(37, 99, 235)
GREEN = RGBColor(15, 159, 110)
ORANGE = RGBColor(249, 115, 22)
PURPLE = RGBColor(124, 58, 237)
RED = RGBColor(220, 38, 38)
YELLOW = RGBColor(255, 247, 237)
BLUE_SOFT = RGBColor(239, 246, 255)
GREEN_SOFT = RGBColor(236, 253, 245)
PURPLE_SOFT = RGBColor(245, 243, 255)


LESSONS = [
    {
        "week": 1,
        "title": "AI Teammate Foundations",
        "subtitle": "ChatGPT, Prompting, and First Assistant",
        "color": BLUE,
        "curriculum": "Students learn that AI is a creative teammate, not just a search box. They practice prompt structure, test answers, and build a personal assistant they can improve.",
        "tools": ["ChatGPT", "Prompt Engineering", "Role + Task + Context + Format", "Critique"],
        "goal": "Move students from simply asking questions to directing an AI teammate with clear roles, context, and output formats.",
        "deliverable": "Personal AI assistant with 5 useful example conversations.",
        "hook": "Can AI do my homework?",
        "big_idea": "The student is the director. AI is the teammate.",
        "concepts": ["AI can take a role", "Specific context improves answers", "A good format makes answers easier to use", "Humans still check the result"],
        "demo_title": "Bad prompt vs. good prompt",
        "demo": ["Bad: Help me study.", "Better: Act as a friendly biology tutor. Quiz me on cells one question at a time. If I miss one, explain it simply."],
        "build": ["Choose an assistant type: Study Buddy, Minecraft Coach, Music Advisor, Sports Trainer, Writing Helper, or Game Idea Coach.", "Create the assistant's role, rules, and answer format.", "Test it with 3 real questions.", "Improve the prompt after seeing the answers."],
        "interactive": ["Thumb vote: helpful or lazy?", "Pair test: ask your partner's assistant one question.", "Repair challenge: turn a vague prompt into a clear prompt."],
        "timing": ["Hook + discussion: 10 min", "Prompt formula demo: 20 min", "Assistant build: 40 min", "Pair test + revise: 15 min", "Exit ticket: 5 min"],
        "prompt": "Act as a friendly study coach for a high school student.\nHelp me study [topic]. Ask one question at a time. If I get it wrong, explain simply and let me try again.",
        "homework": "Create an AI assistant with 5 example conversations. Add one rule that makes the assistant more helpful or safer.",
        "ui_tour": ["New chat: start a clean idea instead of mixing topics.", "Model/tool picker: choose the best available mode for the job.", "/ menu: quick access to tools such as image, canvas, search, file analysis, or research when available.", "Attach/upload: give ChatGPT a file, image, or screenshot to inspect.", "Regenerate/edit: improve a prompt instead of starting over."],
        "slash_tools": ["Search: ask ChatGPT to look up current public information.", "Create image: generate or edit a picture from a description.", "Canvas: open a shared writing/code workspace for longer drafts.", "Analyze file: upload a PDF, spreadsheet, image, or notes for help.", "Deep research: ask for a slower, sourced research report when available.", "More tools: account-specific items may appear or disappear."],
        "bad_prompt": "Help me study.",
        "why_bad": ["No role: ChatGPT does not know whether to be a teacher, quizzer, or note writer.", "No topic: 'study' could mean anything.", "No format: the answer may become a long wall of text.", "No interaction rule: it may answer instead of teaching."],
        "better_prompt": "Act as a friendly biology tutor. Quiz me on cell parts one question at a time. If I miss one, explain it simply and give me another try.",
        "fix_prompt": "Make my prompt better. Ask me for missing details first, then rewrite it using Role + Task + Context + Format.",
        "examples": ["Tutor: Act as an algebra coach. Give me 5 practice problems on slope and explain mistakes.", "Coach: Act as a basketball trainer. Make a 20-minute beginner dribbling plan.", "Creative: Act as a game designer. Give me 3 simple game ideas I can build in HTML."],
        "trivia": {
            "question": "Which prompt will usually teach better?",
            "choices": ["A. Help me study.", "B. Act as a biology tutor. Quiz me on cells one question at a time.", "C. Tell me everything about science."],
            "answer": "B. It gives a role, topic, task, format, and interaction rule.",
        },
    },
    {
        "week": 2,
        "title": "Codex Build Studio",
        "subtitle": "VS Code, Web UI, and Space Invaders Basics",
        "color": GREEN,
        "curriculum": "Students set up the coding workspace, learn files and previews, build a simple web UI, then use Codex to create the basics of a Space Invaders-style browser game.",
        "tools": ["VS Code", "Codex", "HTML", "CSS", "JavaScript", "Browser Preview"],
        "goal": "Help students get a working coding setup and experience the loop of ask, run, debug, and improve.",
        "deliverable": "Personal web page plus a playable Space Invaders starter.",
        "hook": "Can AI become your coding partner inside VS Code?",
        "big_idea": "A project is a folder of files that you can run, inspect, and improve one change at a time.",
        "concepts": ["Open a project folder", "Create and edit files", "Preview a web page", "Describe bugs clearly", "Keep games tiny before adding features"],
        "demo_title": "From web page to tiny game",
        "demo": ["Open VS Code and a course folder.", "Ask Codex to create index.html.", "Preview the page.", "Ask Codex to explain the file.", "Create a tiny canvas game loop."],
        "build": ["Create an AI Builder profile page.", "Customize one section.", "Create space-invaders.html.", "Add player movement, bullets, enemies, score, and game over.", "Playtest and fix one bug."],
        "interactive": ["Setup checkpoint: everyone shows the project folder.", "Choice board: pick one page section to customize.", "Final build sprint: make a basic Space Invaders game with Codex help."],
        "timing": ["Setup + folder tour: 15 min", "Profile page demo: 20 min", "Student customization: 20 min", "Space Invaders starter: 30 min", "Playtest + exit ticket: 5 min"],
        "prompt": "Create a beginner-friendly Space Invaders-style browser game in one file named space-invaders.html. Use HTML, CSS, and JavaScript. Include a player ship, left/right movement, shooting, falling enemies, score, lives, restart button, and comments explaining the main code.",
        "homework": "Expand the Space Invaders starter: add one new enemy type, power-up, sound effect, level, or visual polish. Write 3 sentences explaining what changed.",
        "ui_tour": ["Explorer: shows files in the project folder.", "Editor: where index.html and space-invaders.html are changed.", "Preview/browser: run the page and test it.", "Terminal: optional place to run simple commands.", "Codex chat: ask for one focused change at a time."],
        "bad_prompt": "Make me a game.",
        "why_bad": ["No game type: Codex must guess the rules.", "No file target: it may create multiple files or the wrong file.", "No beginner limit: the code may become too complicated.", "No must-have features: score, lives, restart, and controls may be missing."],
        "better_prompt": "Create a beginner-friendly Space Invaders-style game in one HTML file. Include player movement, shooting, falling enemies, score, lives, restart, and comments.",
        "fix_prompt": "Rewrite this as a Codex prompt with exact file name, controls, features, and a request to explain the code after creating it.",
        "examples": ["Small fix: The bullets do not hit enemies. Find the collision bug and change only the needed code.", "Upgrade: Add a power-up that lets the player shoot two bullets for 8 seconds.", "Explain: Explain the game loop, player movement, and collision detection in plain English."],
        "trivia": {
            "question": "What is the best first bug report to give Codex?",
            "choices": ["A. It is broken.", "B. Fix everything.", "C. Bullets appear, but they do not disappear when they hit enemies."],
            "answer": "C. It describes what happened, what should happen, and where to look.",
        },
    },
    {
        "week": 3,
        "title": "AI Agents and Research Workflows",
        "subtitle": "MCP, YouTube Research, and Browser Safety",
        "color": PURPLE,
        "curriculum": "Students learn how agents plan steps, use tools, and need human supervision. MCP, YouTube research, browser use, and computer use are combined into one practical workflow lesson.",
        "tools": ["AI Agent", "MCP", "YouTube Research", "Browser Use", "Verification", "Safety Rules"],
        "goal": "Help students understand tool-using AI by designing a safe research workflow with clear human checkpoints.",
        "deliverable": "Top 3 video comparison report plus a safe browser task design.",
        "hook": "Can AI use internet tools by itself?",
        "big_idea": "An agent is AI plus a job plus tools.",
        "concepts": ["Chatbots answer inside a conversation", "Agents plan steps and use tools", "MCP connects AI to external tools", "Research needs verification", "Browser tasks need boundaries"],
        "demo_title": "Research agent workflow",
        "demo": ["Choose a student-selected topic.", "Find or simulate three video results.", "Summarize each result.", "Compare which video is best for beginners.", "Mark what a human must verify."],
        "build": ["Pick a topic.", "Create a Top 3 video comparison.", "Write why each video is useful.", "Design a five-step browser task.", "Add one human approval checkpoint."],
        "interactive": ["Sort cards: chatbot or agent?", "Red/yellow/green: is this browser task safe?", "Verification checkpoint: what should humans check?"],
        "timing": ["Agent vs chatbot hook: 10 min", "MCP + tool demo: 20 min", "YouTube research build: 30 min", "Browser safety design: 20 min", "Share-out: 10 min"],
        "prompt": "Act as a careful research assistant. Topic: [student topic]. Compare three beginner-friendly YouTube videos. Summarize each, recommend the best first video, and list what a human should verify before trusting the report.",
        "homework": "Create a Top 3 video comparison report and add a safe 5-step browser task that could extend the research.",
        "ui_tour": ["Chatbot mode: good for questions and summaries.", "Agent workflow: goal, plan, tool use, check result, report back.", "MCP idea: a connector that lets AI use approved tools.", "Browser task: public, supervised, and reversible.", "Verification: humans check facts, links, dates, and source quality."],
        "bad_prompt": "Find videos about coding.",
        "why_bad": ["No audience: beginner, advanced, parent, or teacher?", "No comparison rule: it may list random videos.", "No trust check: it may skip verification.", "No output format: the report may be hard to grade."],
        "better_prompt": "Act as a careful research assistant. Compare 3 beginner-friendly YouTube videos about making a simple JavaScript game. For each, give topic, why it helps beginners, possible weakness, and which one to watch first.",
        "fix_prompt": "Turn this into a safe research-agent prompt. Add source checking, beginner criteria, and a table format.",
        "examples": ["Verification: What facts should I check before trusting this video summary?", "Browser task: Design a safe 5-step public-information task with one human approval checkpoint.", "Report: Make a table with video title, skill level, strongest part, weakness, and recommendation."],
        "trivia": {
            "question": "Which agent task is safest for class?",
            "choices": ["A. Buy the cheapest laptop online.", "B. Compare 3 public tutorials and cite what to verify.", "C. Log into my email and summarize messages."],
            "answer": "B. It uses public information and keeps a human verification step.",
        },
    },
    {
        "week": 4,
        "title": "Final Agent Project",
        "subtitle": "OpenClaw, OAuth, Demo, and Next Build",
        "color": BLUE,
        "curriculum": "Students connect prompting, coding, games, agents, MCP, browser safety, OpenClaw, and OAuth into a final project concept or demo they can explain.",
        "tools": ["OpenClaw", "OAuth", "Codex OAuth", "AI Agent Workflow", "Presentation", "Peer Feedback"],
        "goal": "Help students turn the course pieces into a final project demo, workflow diagram, or next-build plan.",
        "deliverable": "Final AI project concept, demo, or 2-minute presentation.",
        "hook": "Can we connect AI tools like apps on a phone?",
        "big_idea": "A strong AI project is useful, safe, testable, and explainable.",
        "concepts": ["Agent frameworks organize workflows", "OAuth grants limited permission", "Connected tools need trust and safety", "Final projects should be explainable", "Good demos show one working moment"],
        "demo_title": "Agent workflow concept",
        "demo": ["Show a simple OpenClaw-style workflow.", "Explain where login/permission fits.", "Connect the idea back to Codex and tools.", "Model a 2-minute final presentation.", "Show how to ask AI for a next-build checklist."],
        "build": ["Choose final project idea.", "Define user, problem, tools, and workflow.", "Prepare a short demo or concept pitch.", "Give peer feedback.", "Write the next 3 improvements."],
        "interactive": ["App connection analogy: what permissions would this need?", "Gallery walk: leave one feedback note.", "Next-build vote: what should this project add next?"],
        "timing": ["OpenClaw + OAuth concept: 20 min", "Project planning: 25 min", "Demo/pitch build: 25 min", "Gallery walk: 15 min", "Reflection: 5 min"],
        "prompt": "Help me prepare a 2-minute project presentation. Include: what I built, who it helps, best prompt, what broke, how I fixed it, and what I would add next.",
        "homework": "Refine the final project after feedback. Add one improvement to the game, assistant, research workflow, or agent plan and prepare a short demo.",
        "ui_tour": ["Project board: user, problem, tool, input, output, safety check.", "Demo path: one action that works clearly.", "Feedback notes: one strength, one question, one next step.", "OAuth idea: permission without sharing passwords directly.", "Final prompt: ask AI for a tighter demo script and next-build checklist."],
        "bad_prompt": "Make my presentation good.",
        "why_bad": ["No project details: AI cannot tell what matters.", "No audience: a parent pitch differs from a technical demo.", "No time limit: the response may be too long.", "No structure: it may miss what broke and what changed."],
        "better_prompt": "Help me write a 2-minute demo script for my Space Invaders project. Audience: classmates. Include what I built, one bug I fixed, one prompt that helped, and one feature I would add next.",
        "fix_prompt": "Ask me 5 questions about my project, then turn my answers into a 2-minute presentation script.",
        "examples": ["Assistant project: Show one before/after answer from your AI assistant.", "Game project: Show the game running, one bug fix, and one new feature.", "Agent project: Show the workflow diagram and explain where a human approves the action."],
        "trivia": {
            "question": "What makes a final demo strong?",
            "choices": ["A. It shows one clear working moment.", "B. It lists every idea you ever had.", "C. It hides what broke."],
            "answer": "A. A good demo is useful, safe, testable, and easy to explain.",
        },
    },
]


def fill_shape(shape, color, line=None):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = line or color


def fill_background(slide, color):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = color


def rect(slide, x, y, w, h, color, line=None, radius=True):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    fill_shape(shape, color, line)
    return shape


def text(slide, x, y, w, h, value, size=18, bold=False, color=INK, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = value
    if align:
        p.alignment = align
    run = p.runs[0]
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def bullets(slide, x, y, w, h, items, size=22, color=INK):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(9)
    return box


def footer(slide, item, slide_no, total):
    text(slide, 0.45, 7.13, 5.8, 0.18, f"AI Builders Academy | Week {item['week']}", 8.5, False, MUTED)
    text(slide, 12.0, 7.13, 0.9, 0.18, f"{slide_no}/{total}", 8.5, False, MUTED, PP_ALIGN.RIGHT)


def title_bar(slide, item):
    rect(slide, 0, 0, 13.333, 0.72, item["color"], item["color"], radius=False)
    text(slide, 0.55, 0.18, 1.25, 0.22, f"WEEK {item['week']}", 11, True, WHITE)
    text(slide, 1.72, 0.12, 7.1, 0.35, item["title"], 20, True, WHITE)
    text(slide, 8.65, 0.19, 4.1, 0.2, item["subtitle"], 10.5, False, WHITE, PP_ALIGN.RIGHT)


def card(slide, x, y, w, h, title, body, accent=BLUE, bg=WHITE):
    rect(slide, x, y, w, h, bg, LINE)
    rect(slide, x, y, 0.08, h, accent, accent, radius=False)
    text(slide, x + 0.18, y + 0.14, w - 0.35, 0.24, title, 13, True, accent)
    if isinstance(body, list):
        bullets(slide, x + 0.22, y + 0.48, w - 0.42, h - 0.55, body, 14)
    else:
        text(slide, x + 0.18, y + 0.48, w - 0.36, h - 0.55, body, 14)


def arrow(slide, x, y, w, h, label, color=ORANGE):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    fill_shape(shape, color, color)
    shape.text_frame.text = label
    p = shape.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.name = "Aptos"
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = WHITE
    return shape


def prompt_box(slide, value):
    rect(slide, 0.9, 2.05, 11.55, 3.05, DARK, DARK)
    text(slide, 1.18, 2.32, 11.0, 2.45, value, 21, False, WHITE)


def screenshot_box(slide, x, y, w, h, title, rows, accent=BLUE):
    rect(slide, x, y, w, h, WHITE, LINE)
    rect(slide, x, y, w, 0.42, RGBColor(241, 245, 249), LINE, radius=False)
    text(slide, x + 0.18, y + 0.12, w - 0.35, 0.14, title, 8.5, True, MUTED)
    rect(slide, x + 0.28, y + 0.7, w - 0.56, 0.5, RGBColor(248, 250, 252), LINE)
    text(slide, x + 0.46, y + 0.86, w - 0.92, 0.12, "/  Type a tool or command", 9.5, False, MUTED)
    row_y = y + 1.42
    for row in rows[:6]:
        rect(slide, x + 0.28, row_y, w - 0.56, 0.44, RGBColor(255, 255, 255), LINE)
        text(slide, x + 0.46, row_y + 0.13, 0.28, 0.11, "/", 9, True, accent)
        text(slide, x + 0.78, row_y + 0.11, w - 1.25, 0.15, row, 8.7, False, INK)
        row_y += 0.5


def download_path_slide(prs, item, slide_no, total):
    slide = add_slide(prs, item, slide_no, total, "Download everything before class", "Use this when students need the full course folder.")
    card(slide, 0.65, 1.9, 2.25, 2.55, "1. Open repo", "Go to the GitHub page for the class materials.", item["color"], WHITE)
    arrow(slide, 2.95, 2.65, 0.95, 0.45, "next", ORANGE)
    card(slide, 4.0, 1.9, 2.25, 2.55, "2. Code button", "Click the green Code button near the file list.", GREEN, GREEN_SOFT)
    arrow(slide, 6.3, 2.65, 0.95, 0.45, "next", ORANGE)
    card(slide, 7.35, 1.9, 2.25, 2.55, "3. Download ZIP", "Choose Download ZIP from the menu.", BLUE, BLUE_SOFT)
    arrow(slide, 9.65, 2.65, 0.95, 0.45, "next", ORANGE)
    card(slide, 10.7, 1.9, 2.0, 2.55, "4. Open folder", "Unzip it, then open the folder in VS Code.", PURPLE, PURPLE_SOFT)
    card(slide, 1.0, 5.15, 11.2, 0.85, "Teacher note", "If the GitHub layout changes, use the same idea: find the repo, open the download menu, download ZIP, unzip, then open the folder.", ORANGE, YELLOW)
    return slide


def github_screenshot_slide(prs, item, slide_no, total):
    slide = add_slide(prs, item, slide_no, total, "Screenshot guide: GitHub download", "Mock screenshot with arrows for the student walkthrough.")
    rect(slide, 0.85, 1.82, 11.65, 4.75, WHITE, LINE)
    rect(slide, 0.85, 1.82, 11.65, 0.52, RGBColor(248, 250, 252), LINE, radius=False)
    text(slide, 1.1, 1.99, 3.5, 0.16, "github.com / class-repo", 10, True, MUTED)
    text(slide, 1.15, 2.65, 4.8, 0.45, "AI-Tutor", 26, True, INK)
    rect(slide, 8.9, 2.62, 1.55, 0.48, GREEN, GREEN)
    text(slide, 9.15, 2.78, 1.0, 0.13, "Code", 11, True, WHITE, PP_ALIGN.CENTER)
    arrow(slide, 7.2, 2.62, 1.45, 0.45, "click", ORANGE)
    rect(slide, 8.65, 3.25, 2.25, 1.65, WHITE, LINE)
    text(slide, 8.9, 3.48, 1.55, 0.16, "Clone", 9, True, MUTED)
    text(slide, 8.9, 3.93, 1.55, 0.16, "Open with GitHub Desktop", 8.5, False, MUTED)
    rect(slide, 8.82, 4.28, 1.9, 0.38, BLUE_SOFT, LINE)
    text(slide, 9.0, 4.41, 1.55, 0.12, "Download ZIP", 9, True, INK)
    arrow(slide, 6.95, 4.2, 1.6, 0.45, "choose", BLUE)
    card(slide, 1.15, 3.35, 4.95, 1.35, "Say this out loud", ["The Code button is not for writing code today.", "It is the menu where GitHub hides the download option."], item["color"], BLUE_SOFT)
    return slide


def open_folder_screenshot_slide(prs, item, slide_no, total):
    slide = add_slide(prs, item, slide_no, total, "Screenshot guide: open the folder", "After download, students need to unzip and open the project folder.")
    rect(slide, 0.8, 1.9, 5.45, 4.4, WHITE, LINE)
    rect(slide, 0.8, 1.9, 5.45, 0.45, RGBColor(248, 250, 252), LINE, radius=False)
    text(slide, 1.05, 2.05, 2.6, 0.12, "Downloads", 9, True, MUTED)
    card(slide, 1.25, 2.65, 4.15, 0.62, "AI-Tutor-main.zip", "Double-click to unzip.", BLUE, BLUE_SOFT)
    arrow(slide, 5.6, 2.78, 1.1, 0.45, "unzip", ORANGE)
    card(slide, 1.25, 3.65, 4.15, 0.62, "AI-Tutor-main", "This is the folder to open.", GREEN, GREEN_SOFT)
    rect(slide, 7.0, 1.9, 5.45, 4.4, WHITE, LINE)
    rect(slide, 7.0, 1.9, 5.45, 0.45, RGBColor(248, 250, 252), LINE, radius=False)
    text(slide, 7.25, 2.05, 2.9, 0.12, "VS Code", 9, True, MUTED)
    text(slide, 7.35, 2.7, 3.6, 0.35, "File > Open Folder...", 24, True, INK)
    arrow(slide, 6.55, 3.25, 1.2, 0.45, "open", BLUE)
    card(slide, 7.35, 3.75, 4.25, 1.0, "Check", "The left Explorer panel should show folders like slides, docs, scripts, or your project files.", PURPLE, PURPLE_SOFT)
    return slide


def trivia_question_slide(prs, item, slide_no, total):
    slide = add_slide(prs, item, slide_no, total, "Trivia check", "Press next to reveal the answer.")
    text(slide, 0.9, 2.0, 11.5, 0.55, item["trivia"]["question"], 34, True, item["color"], PP_ALIGN.CENTER)
    card(slide, 1.1, 3.05, 11.0, 2.15, "Choose one", item["trivia"]["choices"], PURPLE, PURPLE_SOFT)
    card(slide, 1.35, 5.75, 10.5, 0.62, "Animation cue", "Pause here. Let students vote with fingers, cards, or chat. Advance to reveal.", ORANGE, YELLOW)
    return slide


def trivia_answer_slide(prs, item, slide_no, total):
    slide = add_slide(prs, item, slide_no, total, "Trivia answer reveal", "This is a slide-by-slide reveal so it works anywhere.")
    rect(slide, 1.05, 2.15, 11.2, 2.0, GREEN_SOFT, GREEN)
    text(slide, 1.35, 2.72, 10.6, 0.5, item["trivia"]["answer"], 30, True, INK, PP_ALIGN.CENTER)
    card(slide, 1.2, 4.75, 10.9, 0.85, "Connect it back", "Ask: How could we rewrite the weak prompt or weak plan so it includes the missing details?", item["color"], WHITE)
    return slide


def add_slide(prs, item, slide_no, total, title, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_background(slide, LIGHT)
    title_bar(slide, item)
    text(slide, 0.72, 1.02, 11.7, 0.52, title, 31, True, INK)
    if subtitle:
        text(slide, 0.76, 1.62, 11.4, 0.28, subtitle, 15, False, MUTED)
    footer(slide, item, slide_no, total)
    return slide


def build_week_deck(item):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    total = 18

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_background(slide, DARK)
    rect(slide, 8.9, -0.5, 3.2, 3.2, item["color"], item["color"])
    rect(slide, 0.75, 0.78, 1.7, 0.38, item["color"], item["color"])
    text(slide, 0.92, 0.86, 1.35, 0.18, f"WEEK {item['week']}", 11, True, WHITE, PP_ALIGN.CENTER)
    text(slide, 0.75, 1.58, 9.7, 0.82, item["title"], 44, True, WHITE)
    text(slide, 0.78, 2.58, 9.6, 0.5, item["subtitle"], 22, False, RGBColor(226, 232, 240))
    text(slide, 0.8, 4.45, 9.8, 0.7, item["curriculum"], 20, False, RGBColor(226, 232, 240))
    footer(slide, item, 1, total)

    download_path_slide(prs, item, 2, total)
    github_screenshot_slide(prs, item, 3, total)
    open_folder_screenshot_slide(prs, item, 4, total)

    slide = add_slide(prs, item, 5, total, "Today’s target", "Each lesson is planned for about 90 minutes.")
    card(slide, 0.75, 1.92, 3.65, 1.42, "Main goal", item["goal"], item["color"], WHITE)
    card(slide, 4.65, 1.92, 3.55, 1.42, "Student deliverable", item["deliverable"], GREEN, GREEN_SOFT)
    card(slide, 8.45, 1.92, 3.95, 1.42, "Tools today", item["tools"], PURPLE, PURPLE_SOFT)
    card(slide, 1.0, 3.85, 11.25, 1.8, "90-minute flow", item["timing"], ORANGE, YELLOW)

    slide = add_slide(prs, item, 6, total, "Hook", "Start with a question that makes students curious immediately.")
    text(slide, 0.95, 2.05, 11.4, 1.05, item["hook"], 48, True, item["color"], PP_ALIGN.CENTER)
    bullets(slide, 1.35, 4.0, 10.6, 1.35, ["Ask for quick predictions.", "Let students argue both sides.", "Then show the demo before explaining too much."], 23)

    trivia_question_slide(prs, item, 7, total)
    trivia_answer_slide(prs, item, 8, total)

    slide = add_slide(prs, item, 9, total, "What to show first", "Give students a visible map before asking them to build.")
    if item.get("slash_tools"):
        screenshot_box(slide, 0.8, 1.95, 5.7, 3.8, "Screenshot-style mockup: ChatGPT / menu", item["slash_tools"], item["color"])
        card(slide, 6.85, 1.95, 5.45, 3.8, "Explain it simply", ["Typing / is a quick way to look for available tools.", "The exact list can change by account, plan, and school settings.", "Students do not need every tool today; they need to know how to choose the right one.", "Rule: pick the tool that matches the job, then give clear instructions."], item["color"], BLUE_SOFT)
    else:
        card(slide, 0.85, 1.95, 5.45, 3.8, "UI/tool tour", item["ui_tour"], item["color"], WHITE)
        card(slide, 6.75, 1.95, 5.45, 3.8, "Tutor script", ["Point to each tool before using it.", "Say what it is for in one sentence.", "Use one tiny example before the full activity.", "Ask: Which tool should we use for this job?"], ORANGE, YELLOW)

    slide = add_slide(prs, item, 10, total, "Bad prompt, better prompt, repair prompt", "Students need to see why a prompt fails and how to fix it.")
    card(slide, 0.7, 1.95, 3.55, 3.7, "Bad prompt", item["bad_prompt"], RED, WHITE)
    card(slide, 4.45, 1.95, 3.95, 3.7, "Why it is weak", item["why_bad"], ORANGE, YELLOW)
    card(slide, 8.6, 1.95, 3.95, 1.72, "Better prompt", item["better_prompt"], GREEN, GREEN_SOFT)
    card(slide, 8.6, 3.95, 3.95, 1.7, "Repair prompt", item["fix_prompt"], item["color"], BLUE_SOFT)

    slide = add_slide(prs, item, 11, total, "Examples students can copy", "Use examples as starting points, not scripts to memorize.")
    card(slide, 0.95, 1.95, 11.35, 3.35, "Example prompts", item["examples"], item["color"], WHITE)
    card(slide, 1.0, 5.58, 11.2, 0.78, "How to adapt", "Replace the topic, audience, output format, and difficulty level. Then test once and ask for one improvement.", GREEN, GREEN_SOFT)

    slide = add_slide(prs, item, 12, total, "Big idea", "One concept students should remember after class.")
    text(slide, 1.0, 2.0, 11.3, 0.75, item["big_idea"], 34, True, INK, PP_ALIGN.CENTER)
    card(slide, 1.1, 3.35, 11.1, 1.8, "Key concepts", item["concepts"], item["color"], WHITE)

    slide = add_slide(prs, item, 13, total, item["demo_title"], "Live demo: students see the result before the explanation.")
    card(slide, 0.95, 2.0, 5.25, 3.4, "Demo steps", item["demo"], item["color"], WHITE)
    card(slide, 6.55, 2.0, 5.75, 3.4, "Tutor move", ["Narrate what you are asking AI to do.", "Pause after the output.", "Ask: What should we improve next?"], ORANGE, YELLOW)

    slide = add_slide(prs, item, 14, total, "Live prompt", "Copy, modify, and use during the demo.")
    prompt_box(slide, item["prompt"])
    text(slide, 1.05, 5.55, 10.8, 0.35, "After the response, ask students: What worked? What is missing? What should we ask next?", 16, True, item["color"], PP_ALIGN.CENTER)

    slide = add_slide(prs, item, 15, total, "Hands-on build", "Students create their own version and improve through iteration.")
    card(slide, 0.95, 2.0, 5.45, 3.55, "Build steps", item["build"], GREEN, GREEN_SOFT)
    card(slide, 6.75, 2.0, 5.45, 3.55, "Keep it manageable", ["Start with the smallest working version.", "Let students customize one thing.", "Have them explain the change in plain English."], item["color"], WHITE)

    slide = add_slide(prs, item, 16, total, "Mini challenge", "A short activity that acts like an animation break.")
    card(slide, 0.95, 2.0, 5.45, 3.25, "Challenge", ["Pick one example prompt.", "Change the audience, topic, or feature.", "Run or discuss the new version.", "Share what changed in one sentence."], item["color"], WHITE)
    card(slide, 6.75, 2.0, 5.45, 3.25, "Fast finishers", ["Add one constraint.", "Ask AI to critique the result.", "Ask for one simpler version.", "Prepare a 20-second explanation."], ORANGE, YELLOW)

    slide = add_slide(prs, item, 17, total, "Interactive checkpoints", "Use these to keep kids active instead of just watching.")
    card(slide, 1.0, 2.0, 11.2, 2.3, "Student interaction", item["interactive"], PURPLE, PURPLE_SOFT)
    card(slide, 1.0, 4.75, 11.2, 0.95, "Exit ticket", "What did you make? What prompt helped? What would you change next?", GREEN, GREEN_SOFT)

    slide = add_slide(prs, item, 18, total, "Homework", "Editable: change this based on what you want students to do.")
    rect(slide, 0.95, 2.0, 11.45, 2.0, YELLOW, RGBColor(253, 186, 116))
    text(slide, 1.25, 2.35, 10.8, 0.55, item["homework"], 25, True, INK, PP_ALIGN.CENTER)
    card(slide, 1.0, 4.65, 11.2, 1.0, "Assessment reminder", "Assessment is based on projects, not tests. Students should be able to show what they made and explain one improvement.", ORANGE, WHITE)

    path = OUT_DIR / f"Week_{item['week']}_AI_Builders_{item['title'].replace(' ', '_').replace('&', 'and').replace('+', 'and')}.pptx"
    prs.save(path)
    return path


def main():
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    for item in LESSONS:
        print(build_week_deck(item))


if __name__ == "__main__":
    main()
