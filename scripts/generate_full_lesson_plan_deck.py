from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


OUT = "slides/AI_Builders_Academy_Full_Lesson_Plan_Google_Slides.pptx"

INK = RGBColor(23, 32, 51)
MUTED = RGBColor(86, 101, 123)
WHITE = RGBColor(255, 255, 255)
LIGHT = RGBColor(244, 247, 251)
LINE = RGBColor(217, 226, 238)
DARK = RGBColor(15, 23, 42)
BLUE = RGBColor(37, 99, 235)
GREEN = RGBColor(15, 159, 110)
ORANGE = RGBColor(249, 115, 22)
PURPLE = RGBColor(124, 58, 237)
RED = RGBColor(220, 38, 38)
YELLOW_BG = RGBColor(255, 247, 237)
BLUE_BG = RGBColor(239, 246, 255)
GREEN_BG = RGBColor(236, 253, 245)
PURPLE_BG = RGBColor(245, 243, 255)


def fill(shape, color, line=None):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = line or color


def fill_background(slide, color):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = color


def textbox(slide, x, y, w, h, text, size=18, bold=False, color=INK, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = text
    if align:
        p.alignment = align
    r = p.runs[0]
    r.font.name = "Aptos"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return box


def bullets(slide, x, y, w, h, items, size=15, color=INK):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(5)
    return box


def rect(slide, x, y, w, h, color, line=None, radius=True):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    fill(shape, color, line)
    return shape


def card(slide, x, y, w, h, title, body=None, accent=BLUE, bg=WHITE):
    rect(slide, x, y, w, h, bg, LINE)
    rect(slide, x, y, 0.08, h, accent, accent, radius=False)
    textbox(slide, x + 0.18, y + 0.15, w - 0.34, 0.26, title, 14, True, accent)
    if body:
        textbox(slide, x + 0.18, y + 0.48, w - 0.34, h - 0.55, body, 12.5, False, INK)


def prompt_box(slide, x, y, w, h, text):
    rect(slide, x, y, w, h, DARK, DARK)
    textbox(slide, x + 0.18, y + 0.18, w - 0.36, h - 0.25, text, 11.5, False, WHITE)


def homework_box(slide, x, y, w, h, text):
    rect(slide, x, y, w, h, YELLOW_BG, RGBColor(253, 186, 116))
    textbox(slide, x + 0.18, y + 0.12, w - 0.36, 0.24, "Homework - editable", 13, True, ORANGE)
    textbox(slide, x + 0.18, y + 0.45, w - 0.36, h - 0.48, text, 12.5, False, INK)


def title_bar(slide, week, title, subtitle, color):
    rect(slide, 0, 0, 13.333, 0.68, color, color, radius=False)
    textbox(slide, 0.55, 0.17, 1.25, 0.24, f"WEEK {week}", 11, True, WHITE)
    textbox(slide, 1.65, 0.12, 6.9, 0.36, title, 20, True, WHITE)
    textbox(slide, 8.55, 0.18, 4.2, 0.24, subtitle, 10.5, False, WHITE, PP_ALIGN.RIGHT)


def footer(slide, n):
    textbox(slide, 0.45, 7.15, 6.5, 0.18, "AI Builders Academy | Full tutor lesson plan", 8.5, False, MUTED)
    textbox(slide, 12.5, 7.15, 0.45, 0.18, str(n), 8.5, False, MUTED, PP_ALIGN.RIGHT)


def nav_button(slide, x, label, target=None):
    shape = rect(slide, x, 6.85, 1.25, 0.32, WHITE, LINE)
    shape.text_frame.text = label
    p = shape.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.runs[0].font.name = "Aptos"
    p.runs[0].font.size = Pt(8.5)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = BLUE
    if target is not None:
        shape.click_action.target_slide = target
    return shape


def background(slide):
    fill_background(slide, LIGHT)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # Title
    s = prs.slides.add_slide(blank)
    fill_background(s, DARK)
    rect(s, 8.85, -0.55, 3.5, 3.5, ORANGE, ORANGE)
    rect(s, 10.25, 4.8, 2.15, 2.15, GREEN, GREEN)
    rect(s, 0.75, 0.75, 2.35, 0.38, BLUE, BLUE)
    textbox(s, 0.9, 0.83, 2.05, 0.18, "Tutor Lesson Plan", 11, True, WHITE, PP_ALIGN.CENTER)
    textbox(s, 0.75, 1.55, 8.8, 0.85, "AI Builders Academy", 44, True, WHITE)
    textbox(s, 0.78, 2.55, 9.5, 0.95, "Four-lesson vibe coding curriculum for tutors teaching kids to build with ChatGPT, VS Code, Codex, games, agents, and safe AI workflows.", 22, False, RGBColor(226, 232, 240))
    textbox(s, 0.78, 5.72, 8.2, 0.32, "Use the Lesson Hub for direct links to each 90-minute lesson.", 14, False, RGBColor(203, 213, 225))
    footer(s, 1)
    title = s

    # Hub
    hub = prs.slides.add_slide(blank)
    fill_background(hub, LIGHT)
    textbox(hub, 0.55, 0.45, 8.6, 0.52, "Lesson Hub: four direct lesson links", 30, True, INK)
    textbox(hub, 0.58, 1.02, 9.7, 0.25, "Click a lesson card during tutor prep or while teaching. Each lesson is designed as one complete 90-minute plan.", 12.5, False, MUTED)
    footer(hub, 2)

    weeks = []
    week_data = [
        {
            "week": 1,
            "title": "AI Teammate Foundations",
            "subtitle": "ChatGPT, prompting, and first assistant",
            "color": BLUE,
            "goal": "Students learn to direct AI with clear roles, context, and output formats.",
            "output": "Personal AI assistant with 5 example conversations.",
            "flow": ["0-10 Hook: Can AI do my homework?", "10-30 Prompt formula demo", "30-70 Build custom assistant", "70-85 Partner test + revise", "85-90 Exit ticket"],
            "interactive": ["Thumb vote: helpful or lazy?", "Pair test: ask your partner's assistant one question", "Debug moment: make a vague prompt clearer"],
            "prompt": "Act as a friendly study coach for a high school student.\nHelp me study [topic]. Ask one question at a time. If I get it wrong, explain simply and let me try again.",
            "homework": "Create 5 example conversations with your assistant and add one safety rule.",
        },
        {
            "week": 2,
            "title": "Codex Build Studio",
            "subtitle": "VS Code, web UI, and Space Invaders basics",
            "color": GREEN,
            "goal": "Students set up Codex, create a web UI, and build the basics of a browser game.",
            "output": "Personal web page plus a playable Space Invaders starter.",
            "flow": ["0-15 Setup + folder tour", "15-35 Profile page demo", "35-55 Student customization", "55-85 Space Invaders starter", "85-90 Playtest + exit"],
            "interactive": ["Setup checkpoint: everyone shows folder", "Choice board: customize one section", "Final sprint: build Space Invaders basics"],
            "prompt": "Create a beginner-friendly Space Invaders-style browser game in one file named space-invaders.html. Include a player ship, left/right movement, shooting, falling enemies, score, lives, restart button, and comments explaining the main code.",
            "homework": "Expand the Space Invaders starter with one enemy type, power-up, sound effect, level, or visual polish.",
        },
        {
            "week": 3,
            "title": "AI Agents and Research Workflows",
            "subtitle": "MCP, YouTube research, and browser safety",
            "color": PURPLE,
            "goal": "Students design a safe tool-using research workflow with human verification.",
            "output": "Top 3 video report plus safe browser task design.",
            "flow": ["0-10 Agent vs chatbot hook", "10-30 MCP + tools demo", "30-60 YouTube research build", "60-80 Browser safety design", "80-90 Share-out"],
            "interactive": ["Sort cards: chatbot or agent?", "Red/yellow/green: is this task safe?", "Verification checkpoint: what should humans check?"],
            "prompt": "Act as a careful research assistant. Topic: [student topic]. Compare three beginner-friendly YouTube videos. Summarize each, recommend the best first video, and list what a human should verify.",
            "homework": "Create a Top 3 video report and add a safe 5-step browser task that could extend the research.",
        },
        {
            "week": 4,
            "title": "Final Agent Project",
            "subtitle": "OpenClaw, OAuth, demo, and next build",
            "color": BLUE,
            "goal": "Students connect course ideas into a final project concept, demo, or workflow presentation.",
            "output": "Final project concept, demo, or 2-minute presentation.",
            "flow": ["0-20 OpenClaw + OAuth concept", "20-45 Project planning", "45-70 Demo/pitch build", "70-85 Gallery walk", "85-90 Reflection"],
            "interactive": ["Gallery walk: leave one note", "Presentation timer: 2 minutes each", "Next-build vote: what would you add?"],
            "prompt": "Help me prepare a 2-minute project presentation. Include: what I built, who it helps, best prompt, what broke, how I fixed it, and what I would add next.",
            "homework": "Refine the final project after feedback and add one improvement to the game, assistant, workflow, or agent plan.",
        },
    ]

    # Create week slides first so hub can link to them.
    for item in week_data:
        s = prs.slides.add_slide(blank)
        fill_background(s, LIGHT)
        title_bar(s, item["week"], item["title"], item["subtitle"], item["color"])
        footer(s, item["week"] + 2)

        card(s, 0.45, 0.92, 3.0, 1.15, "Goal", item["goal"], item["color"], WHITE)
        card(s, 3.65, 0.92, 2.55, 1.15, "Student output", item["output"], GREEN, WHITE)
        homework_box(s, 6.4, 0.92, 6.45, 1.15, item["homework"])

        card(s, 0.45, 2.32, 3.85, 2.0, "90-minute flow", None, item["color"], WHITE)
        bullets(s, 0.72, 2.85, 3.35, 1.15, item["flow"], 11.5)

        card(s, 4.55, 2.32, 3.8, 2.0, "Interactive moments", None, PURPLE, PURPLE_BG)
        bullets(s, 4.82, 2.85, 3.25, 1.05, item["interactive"], 11.5)

        card(s, 8.6, 2.32, 4.25, 2.0, "Tutor moves", None, ORANGE, BLUE_BG)
        bullets(s, 8.87, 2.85, 3.72, 1.08, [
            "Ask students to predict before AI answers.",
            "Pause after every demo to name what changed.",
            "Make one small improvement, then test."
        ], 11.2)

        textbox(s, 0.55, 4.62, 2.2, 0.25, "Live demo prompt", 15, True, DARK)
        prompt_box(s, 0.45, 4.95, 8.0, 1.5, item["prompt"])

        card(s, 8.72, 4.95, 4.13, 1.5, "Exit ticket", "Before students leave, ask: What did you make? What was one prompt that helped? What is one thing you would change next?", GREEN, GREEN_BG)

        weeks.append(s)

    # Hub cards after week slides exist.
    coords = [(0.65, 1.65), (4.72, 1.65), (8.79, 1.65), (4.72, 4.05)]
    for i, item in enumerate(week_data):
        x, y = coords[i]
        shape = rect(hub, x, y, 3.55, 1.65, WHITE, LINE)
        shape.click_action.target_slide = weeks[i]
        rect(hub, x, y, 0.12, 1.65, item["color"], item["color"], radius=False)
        textbox(hub, x + 0.28, y + 0.2, 2.9, 0.22, f"Lesson {item['week']}", 12, True, item["color"])
        textbox(hub, x + 0.28, y + 0.48, 2.95, 0.38, item["title"], 16, True, INK)
        textbox(hub, x + 0.28, y + 0.96, 2.9, 0.35, item["output"], 10.5, False, MUTED)
        textbox(hub, x + 2.82, y + 1.32, 0.45, 0.16, "Open", 8, True, item["color"], PP_ALIGN.RIGHT)

    # Add title slide button after hub exists.
    btn = rect(title, 0.78, 4.85, 2.05, 0.48, BLUE, BLUE)
    btn.click_action.target_slide = hub
    btn.text_frame.text = "Open Lesson Hub"
    btn.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    btn.text_frame.paragraphs[0].runs[0].font.name = "Aptos"
    btn.text_frame.paragraphs[0].runs[0].font.size = Pt(13)
    btn.text_frame.paragraphs[0].runs[0].font.bold = True
    btn.text_frame.paragraphs[0].runs[0].font.color.rgb = WHITE

    # Week navigation buttons.
    for i, s in enumerate(weeks):
        nav_button(s, 8.05, "Lesson Hub", hub)
        if i > 0:
            nav_button(s, 9.45, "Prev", weeks[i - 1])
        if i < len(weeks) - 1:
            nav_button(s, 10.85, "Next", weeks[i + 1])

    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    build()
