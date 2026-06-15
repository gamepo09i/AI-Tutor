from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUT = "slides/AI_Builders_Academy_Vibe_Coding_Tutor_Deck.pptx"

COLORS = {
    "ink": RGBColor(23, 32, 51),
    "muted": RGBColor(91, 106, 128),
    "blue": RGBColor(37, 99, 235),
    "green": RGBColor(15, 159, 110),
    "orange": RGBColor(249, 115, 22),
    "purple": RGBColor(124, 58, 237),
    "light": RGBColor(244, 247, 251),
    "white": RGBColor(255, 255, 255),
    "line": RGBColor(220, 228, 240),
    "dark": RGBColor(16, 24, 39),
}


def set_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color


def add_textbox(slide, x, y, w, h, text, size=24, bold=False, color=None, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    if align:
        p.alignment = align
    run = p.runs[0]
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color or COLORS["ink"]
    run.font.name = "Aptos"
    return box


def add_bullets(slide, x, y, w, h, items, size=24, color=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = color or COLORS["ink"]
        p.font.name = "Aptos"
        p.space_after = Pt(8)
    return box


def add_title(slide, title, subtitle=None):
    add_textbox(slide, 0.65, 0.45, 12.0, 0.8, title, size=34, bold=True)
    if subtitle:
        add_textbox(slide, 0.68, 1.18, 11.7, 0.45, subtitle, size=15, color=COLORS["muted"])


def add_footer(slide, number):
    add_textbox(slide, 0.65, 7.12, 5.4, 0.28, "AI Builders Academy | Tutor Deck", size=9, color=COLORS["muted"])
    add_textbox(slide, 12.15, 7.12, 0.55, 0.28, str(number), size=9, color=COLORS["muted"], align=PP_ALIGN.RIGHT)


def add_chip(slide, x, y, text, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(2.25), Inches(0.42))
    set_fill(shape, color)
    shape.text_frame.text = text
    p = shape.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.runs[0].font.size = Pt(12)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = COLORS["white"]


def add_card(slide, x, y, w, h, title, body, color=COLORS["blue"]):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    set_fill(card, COLORS["white"])
    card.line.color.rgb = COLORS["line"]
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(h))
    set_fill(bar, color)
    add_textbox(slide, x + 0.22, y + 0.18, w - 0.45, 0.32, title, size=18, bold=True, color=color)
    add_textbox(slide, x + 0.22, y + 0.62, w - 0.45, h - 0.78, body, size=16, color=COLORS["ink"])


def add_prompt_box(slide, text, y=2.05):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.95), Inches(y), Inches(11.45), Inches(3.25))
    set_fill(box, COLORS["dark"])
    add_textbox(slide, 1.25, y + 0.25, 10.9, 2.75, text, size=20, color=COLORS["white"])


def add_tutor_note(slide, text, y=5.85):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.95), Inches(y), Inches(11.45), Inches(0.78))
    set_fill(shape, RGBColor(255, 247, 237))
    shape.line.color.rgb = RGBColor(253, 186, 116)
    add_textbox(slide, 1.2, y + 0.16, 10.95, 0.43, "Tutor note: " + text, size=14, color=COLORS["ink"])


def blank_slide(prs, number, title=None, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = COLORS["light"]
    if title:
        add_title(slide, title, subtitle)
    add_footer(slide, number)
    return slide


def title_slide(prs, number):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = COLORS["dark"]
    slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(7.5)).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = COLORS["dark"]
    accent = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.2), Inches(-0.6), Inches(3.8), Inches(3.8))
    set_fill(accent, COLORS["orange"])
    accent2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.7), Inches(4.35), Inches(2.2), Inches(2.2))
    set_fill(accent2, COLORS["green"])
    add_chip(slide, 0.85, 0.75, "4-Lesson Tutor Deck", COLORS["blue"])
    add_textbox(slide, 0.85, 1.55, 8.9, 1.75, "AI Builders Academy", size=52, bold=True, color=COLORS["white"])
    add_textbox(slide, 0.9, 3.05, 9.45, 1.05, "Teaching kids vibe coding with ChatGPT, VS Code, Codex, games, agents, and practical AI workflows.", size=24, color=RGBColor(226, 232, 240))
    add_textbox(slide, 0.9, 5.9, 6.2, 0.35, "For tutors: goals, timing, demos, prompts, and assignments", size=15, color=RGBColor(203, 213, 225))
    add_footer(slide, number)


def section_slide(prs, number, title, subtitle, color):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = color
    add_textbox(slide, 0.85, 1.95, 10.7, 1.15, title, size=46, bold=True, color=COLORS["white"])
    add_textbox(slide, 0.9, 3.15, 9.9, 0.85, subtitle, size=23, color=RGBColor(240, 249, 255))
    add_footer(slide, number)
    return slide


def lesson_slide(prs, number, class_no, title, goal, flow, output, homework, color):
    slide = blank_slide(prs, number, f"Lesson {class_no}: {title}", "90-minute teaching plan")
    add_card(slide, 0.78, 1.75, 3.8, 1.45, "Goal", goal, color)
    add_card(slide, 4.78, 1.75, 3.8, 1.45, "Student output", output, COLORS["green"])
    add_card(slide, 8.78, 1.75, 3.8, 1.45, "Homework", homework, COLORS["orange"])
    add_textbox(slide, 0.85, 3.65, 3.0, 0.38, "Suggested flow", size=19, bold=True, color=COLORS["ink"])
    add_bullets(slide, 1.05, 4.08, 11.1, 1.55, flow, size=19)
    add_tutor_note(slide, "Keep the first build small. The win is a working project students can explain.", y=6.15)


def build_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    n = 1
    title_slide(prs, n)

    n += 1
    slide = blank_slide(prs, n, "How to use this deck", "Designed for tutors who may be new to AI coding")
    add_bullets(slide, 0.95, 1.75, 11.45, 3.5, [
        "Use the lesson plan slides to teach the four-session course.",
        "Use the prompt slides as live demo scripts.",
        "Let students build small first, then customize.",
        "Do not worry about knowing every answer. Model asking, testing, and debugging.",
        "Upload this PowerPoint file to Google Drive and open it with Google Slides."
    ], size=24)
    add_tutor_note(slide, "This deck is meant to be edited. Add your school name, dates, and local setup details.")

    n += 1
    slide = blank_slide(prs, n, "What is vibe coding?", "The core teaching idea")
    add_card(slide, 0.88, 1.85, 5.75, 2.2, "Student-friendly definition", "Vibe coding means using AI as a coding teammate. Students describe what they want, inspect what AI makes, test it, and improve it.", COLORS["blue"])
    add_card(slide, 6.88, 1.85, 5.55, 2.2, "Important boundary", "The student is still the director. AI can suggest code, but the student chooses the goal, tests the result, and explains the changes.", COLORS["purple"])
    add_tutor_note(slide, "Say: AI is fast, but it still needs clear instructions and human judgment.")

    n += 1
    slide = blank_slide(prs, n, "Course outcomes", "By the end, students should be able to...")
    add_bullets(slide, 0.95, 1.75, 11.45, 3.7, [
        "Write clearer prompts using role, task, context, and format.",
        "Use VS Code and Codex to create and edit small projects.",
        "Build a simple web UI and a playable browser game.",
        "Explain the difference between a chatbot and an agent.",
        "Design a safe AI workflow using tools, browser tasks, or research steps.",
        "Present what they made, what broke, and how they improved it."
    ], size=23)

    n += 1
    slide = blank_slide(prs, n, "Course roadmap", "Four 90-minute lessons, each with a buildable output")
    roadmap = [
        ("Lesson 1", "AI foundations", "Personal AI assistant"),
        ("Lesson 2", "Codex build studio", "Web UI + Space Invaders starter"),
        ("Lesson 3", "Agents + research", "Video report + browser task"),
        ("Lesson 4", "Final agent project", "Demo or project pitch"),
    ]
    x_positions = [0.75, 4.65, 8.55, 4.65]
    y_positions = [1.75, 1.75, 1.75, 4.25]
    for i, (label, title, body) in enumerate(roadmap):
        add_card(slide, x_positions[i], y_positions[i], 3.55, 1.75, f"{label}: {title}", body, [COLORS["blue"], COLORS["green"], COLORS["purple"], COLORS["blue"]][i])

    n += 1
    slide = blank_slide(prs, n, "Default lesson structure", "Use this rhythm every week")
    add_bullets(slide, 0.95, 1.65, 11.45, 4.05, [
        "Hook, 5-10 min: show a surprising or useful AI example.",
        "Demo, 15-25 min: tutor builds one small thing live.",
        "Hands-on build, 40-60 min: students make their own version.",
        "Share/debug, 10-20 min: students test each other's work.",
        "Wrap-up, 5 min: one takeaway, one assignment, one next-step idea."
    ], size=24)
    add_tutor_note(slide, "If time is short, protect the hands-on build. Students remember what they make.")

    n += 1
    slide = blank_slide(prs, n, "Prompt formula", "A repeatable pattern for students")
    add_card(slide, 0.85, 1.85, 2.75, 2.2, "Role", "Who should AI act like?", COLORS["blue"])
    add_card(slide, 3.85, 1.85, 2.75, 2.2, "Task", "What should AI do?", COLORS["green"])
    add_card(slide, 6.85, 1.85, 2.75, 2.2, "Context", "What details matter?", COLORS["orange"])
    add_card(slide, 9.85, 1.85, 2.75, 2.2, "Format", "How should the answer look?", COLORS["purple"])
    add_prompt_box(slide, "Act as a beginner-friendly coding tutor. Help me build a simple clicker game in one HTML file. Include score, a timer, and comments explaining the code.", y=4.65)

    n += 1
    section_slide(prs, n, "Lesson-by-Lesson Teaching Plan", "Each lesson includes a clear output students can show.", COLORS["blue"])

    lessons = [
        (1, "AI teammate foundations", "Students learn that AI can be a tutor, coach, brainstorming partner, and project helper when prompted clearly.", ["Hook: Can AI do my homework, and should it?", "Demo: bad prompt vs. better prompt.", "Build: students design a custom assistant.", "Share: students test assistant answers with a partner."], "Custom assistant with five example conversations.", "Add five example conversations and one safety rule.", COLORS["blue"]),
        (2, "Codex build studio", "Students set up VS Code/Codex, make a simple web UI, then build a Space Invaders-style starter.", ["Setup: open folder and preview page.", "Demo: create index.html.", "Build: personalize an AI Builder profile page.", "Final sprint: make Space Invaders basics with Codex."], "Personal web page plus playable Space Invaders starter.", "Expand the game with one enemy type, power-up, sound, level, or visual polish.", COLORS["green"]),
        (3, "Agents and research workflows", "Students combine MCP, YouTube research, browser use, verification, and safety boundaries.", ["Hook: Can AI use internet tools by itself?", "Demo: chatbot vs agent.", "Build: top-three YouTube report.", "Safety review: design a browser task with approval checkpoints."], "Video comparison report plus safe browser task design.", "Create a Top 3 video report and a safe 5-step browser task.", COLORS["purple"]),
        (4, "Final agent project", "Students connect the course ideas into a final concept, demo, or workflow presentation.", ["Hook: What could you build next with AI?", "Demo: OpenClaw/OAuth workflow concept.", "Build: prepare final demo or project pitch.", "Present: students explain goal, prompt, bug, and next step."], "Final project concept, demo, or presentation.", "Revise the final project after feedback and add one improvement.", COLORS["blue"]),
    ]
    for lesson in lessons:
        n += 1
        lesson_slide(prs, n, *lesson)

    n += 1
    section_slide(prs, n, "Live Demo Prompts", "Copy these prompts into ChatGPT or Codex during class.", COLORS["green"])

    prompts = [
        ("Lesson 1 demo prompt", "Act as a friendly study coach for a high school student.\n\nHelp me prepare for a biology quiz about cells.\n\nAsk me one question at a time. If I get it wrong, explain it simply and give me another try."),
        ("Lesson 2 web prompt", "Create a beginner-friendly personal web page in one file named index.html.\n\nTheme: AI Builder profile\nSections: name, favorite apps, project ideas, and one button\nStyle: clean, colorful, easy to read\n\nAfter creating it, explain the file in simple language."),
        ("Lesson 2 game prompt", "Create a beginner-friendly Space Invaders-style browser game in one file named space-invaders.html.\n\nInclude player movement, shooting, falling enemies, score, lives, restart, and comments explaining the code."),
        ("Lesson 2 improvement prompt", "Improve this Space Invaders starter for beginners.\n\nAdd one new feature, explain exactly what changed, and keep the code in one HTML file. Do not make it too complicated."),
        ("Lesson 3 research prompt", "Act as a careful research assistant.\n\nTopic: [student topic]\nCompare three useful YouTube video ideas for beginners. Summarize each one, recommend the best first video, and list what a human should verify."),
        ("Lesson 3 task-design prompt", "Design a safe browser task for an AI agent.\n\nThe task must use only public information, must not log into accounts, must not buy anything, and must include a human approval checkpoint before any submission."),
    ]
    for title, prompt in prompts:
        n += 1
        slide = blank_slide(prs, n, title, "Live demo script")
        add_prompt_box(slide, prompt, y=1.85)
        add_tutor_note(slide, "After AI responds, ask students: What worked? What is missing? What should we ask next?", y=5.65)

    n += 1
    section_slide(prs, n, "Tutor Support Slides", "Use these when training other tutors.", COLORS["purple"])

    n += 1
    slide = blank_slide(prs, n, "What tutors should say when stuck", "This is normal in AI classes")
    add_bullets(slide, 0.95, 1.75, 11.45, 3.8, [
        "\"I do not know yet. Let's test it.\"",
        "\"Describe what you expected and what actually happened.\"",
        "\"Ask Codex for the smallest fix, not a total rewrite.\"",
        "\"Before we trust this, how can we verify it?\"",
        "\"Can you explain what changed in your own words?\""
    ], size=25)

    n += 1
    slide = blank_slide(prs, n, "Classroom safety rules", "Especially important for agents and browser tasks")
    add_bullets(slide, 0.95, 1.75, 11.45, 3.9, [
        "Do not type private passwords during demos.",
        "Do not let AI buy anything, submit real forms, or message people without permission.",
        "Use practice sites, public information, or instructor-approved examples.",
        "Watch what the agent is doing and stop it if it goes off-task.",
        "Students should explain the task before running it."
    ], size=24)

    n += 1
    slide = blank_slide(prs, n, "Assessment rubric", "Grade projects, not memorization")
    add_card(slide, 0.75, 1.75, 3.0, 2.1, "Build", "Does the project run or clearly show the intended idea?", COLORS["blue"])
    add_card(slide, 3.95, 1.75, 3.0, 2.1, "Prompting", "Did the student use specific prompts and iterate?", COLORS["green"])
    add_card(slide, 7.15, 1.75, 3.0, 2.1, "Debugging", "Can the student describe one problem and how they fixed it?", COLORS["orange"])
    add_card(slide, 10.35, 1.75, 2.25, 2.1, "Share", "Can they explain the result?", COLORS["purple"])
    add_tutor_note(slide, "A simple 1-4 score for each category is enough.")

    n += 1
    slide = blank_slide(prs, n, "Materials checklist", "Prepare these before class")
    add_bullets(slide, 0.95, 1.75, 11.45, 3.9, [
        "Projector or large display.",
        "Student laptops and internet access.",
        "ChatGPT account path or instructor-provided account plan.",
        "VS Code installed or available through school devices.",
        "Codex access or a fallback pair-programming demo.",
        "Starter folders for web page and game projects.",
        "A backup demo project in case setup takes longer than expected."
    ], size=23)

    n += 1
    slide = blank_slide(prs, n, "Final presentation template", "Give this to students before Lesson 4")
    add_bullets(slide, 0.95, 1.75, 11.45, 3.9, [
        "What did you build?",
        "Who is it for?",
        "What prompt helped the most?",
        "What broke or confused you?",
        "How did you improve it?",
        "What would you add next?"
    ], size=26)

    n += 1
    slide = blank_slide(prs, n, "GitHub sharing workflow", "For the lead tutor")
    add_bullets(slide, 0.95, 1.75, 11.45, 3.7, [
        "Put the useful files in the course folder.",
        "Commit a meaningful batch of changes.",
        "Push to GitHub.",
        "Tell tutors to download the repo or pull the newest version.",
        "Use clear commit messages like: Update lesson 2 game prompts."
    ], size=24)
    add_prompt_box(slide, "git status\ngit add .\ngit commit -m \"Update tutor slide deck\"\ngit push", y=5.15)

    n += 1
    slide = blank_slide(prs, n, "Closing message for tutors", "The tone matters")
    add_card(slide, 0.95, 1.85, 11.35, 2.2, "Main idea", "You do not need to be the person with every answer. Be the person who models how to ask, test, debug, verify, and keep building.", COLORS["green"])
    add_textbox(slide, 1.05, 4.55, 10.8, 0.75, "Build small. Test often. Explain what changed.", size=34, bold=True, color=COLORS["ink"], align=PP_ALIGN.CENTER)

    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    build_deck()
